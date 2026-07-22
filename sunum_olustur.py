import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_arin_ai_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Renk Paleti (Endüstriyel İSG / Maden Teması)
    DARK_BG = RGBColor(15, 23, 42)      # #0F172A (Koyu Lacivert)
    ORANGE = RGBColor(249, 115, 22)     # #F97316 (Arın AI Turuncusu)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(226, 232, 240)
    CARD_BG = RGBColor(30, 41, 59)      # #1E293B

    def add_header(slide, title_text, subtitle_text=""):
        # Başlık
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        
        # Alt Başlık (varsa)
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(14)
            p2.font.color.rgb = LIGHT_GRAY

    # -------------------------------------------------------------
    # SLAYT 1: Kapak
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.color.rgb = DARK_BG

    tb_cover = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf_cover = tb_cover.text_frame
    tf_cover.word_wrap = True

    p_main = tf_cover.paragraphs[0]
    p_main.text = "🛡️ ARIN AI ENTERPRISE"
    p_main.font.size = Pt(44)
    p_main.font.bold = True
    p_main.font.color.rgb = ORANGE
    p_main.alignment = PP_ALIGN.CENTER

    p_sub = tf_cover.add_paragraph()
    p_sub.text = "Maden Sektörüne Özel 360° Proaktif İSG & Karar Destek Platformu"
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = WHITE
    p_sub.alignment = PP_ALIGN.CENTER

    p_author = tf_cover.add_paragraph()
    p_author.text = "\n\nAlperen Taşkıran | Aethel Technologies"
    p_author.font.size = Pt(16)
    p_author.font.color.rgb = LIGHT_GRAY
    p_author.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLAYT 2: Problem (Saha Gerçekleri)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = DARK_BG
    bg2.line.color.rgb = DARK_BG
    add_header(slide2, "📍 Problem: Saha Gerçekleri", "Geleneksel İSG süreçleri reaktiftir ve zaman kaybına yol açar.")

    problems = [
        ("📚 Karmaşık Mevzuat & Veri Yığını", "Yüzlerce sayfalık maden mevzuatı ve geçmiş kaza raporları arasında acil durumlarda hızlı karar vermek oldukça zordur."),
        ("⏳ Vardiyalar Arası Veri Kopukluğu", "Önceki vardiyadaki kılcal gaz artışları veya çatlaklar, sonraki vardiyaya aktarılırken gözden kaçabilmektedir."),
        ("⚡ Operasyonel Görüş Ayrılıkları", "İSG Ekipleri (Tam Güvenlik) ile Üretim Ekipleri (Operasyonel Devamlılık) arasındaki doğal iletişim ve karar çatışmaları.")
    ]

    for i, (title, desc) in enumerate(problems):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ORANGE
        
        tf = card.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = ORANGE
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLAYT 3: Çözüm (Arın AI Nedir?)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = DARK_BG
    bg3.line.color.rgb = DARK_BG
    add_header(slide3, "📍 Çözüm: Arın AI Nedir?", "Saha verilerini anlık analiz eden proaktif karar destek platformu.")

    solutions = [
        ("🔍 360° RAG Mimarisi", "Mevzuat, Tarihsel Kazalar ve MTA Jeoloji veritabanlarını eş zamanlı tarayarak kanıta dayalı analiz sunar."),
        ("🤝 Multi-Agent Çatışma Simülasyonu", "İSG ve Üretim Ajanlarının görüşlerini sentezleyerek Başmühendis Ajanı üzerinden optimum kararı üretir."),
        ("📈 Shift Memory (Saha Hafızası)", "Gaz ivmelenmelerini ve saha risklerini vardiyalar arası kesintisiz olarak takip eder ve görselleştirir.")
    ]

    for i, (title, desc) in enumerate(solutions):
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ORANGE
        
        tf = card.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = ORANGE
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLAYT 4: Öne Çıkan Core Özellikler (Tablo)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = DARK_BG
    bg4.line.color.rgb = DARK_BG
    add_header(slide4, "📍 Öne Çıkan Core Özellikler", "Modüllerin Saha Avantajları ve Teknik Yetenekleri")

    rows, cols = 5, 3
    table_shape = slide4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    table = table_shape.table

    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(4.2)
    table.columns[2].width = Inches(4.3)

    headers = ["Modül / Özellik", "Ne İşe Yarar?", "Saha Avantajı"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_BG

    data = [
        ("📍 Saha Lokasyon Hiyerarşisi", "Ayna, Kör Galeri gibi alanları otomatik tanır.", "Lokasyona özel kritik risk ve kontrol listesi sunar."),
        ("🚜 Ekipman & LOTO Kartı", "Fan, trafo vb. arızalarında otomasyona geçer.", "Zorunlu Kilitleme-Etiketleme prosedürünü basar."),
        ("📈 CH4 Trend Hafızası", "Vardiyalar arası gaz değişimini takip eder.", "Patlama riski gerçekleşmeden proaktif uyarı verir."),
        ("🧮 Deterministik Motorlar", "Fine-Kinney, L-Tipi Matris ve Gürültü hesabı yapar.", "Halüsinasyon riskini sıfırlar, kesin matematik verir.")
    ]

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLAYT 5: İş Akışı (Human-in-the-Loop)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = DARK_BG
    bg5.line.color.rgb = DARK_BG
    add_header(slide5, "📍 Güvenli İş Akışı: Human-in-the-Loop", "Yapay zeka önerir, yetkili başmühendis onaylar.")

    card_workflow = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    card_workflow.fill.solid()
    card_workflow.fill.fore_color.rgb = CARD_BG
    card_workflow.line.color.rgb = ORANGE

    tf_wf = card_workflow.text_frame
    tf_wf.word_wrap = True

    p = tf_wf.paragraphs[0]
    p.text = "🔐 Tam Yasal Uyum ve Sorumluluk Yönetimi\n"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    points = [
        "• Çoklu Girdi Desteği: Ses kayıtları (Whisper), PDF, Word, Excel ve manuel metinler işlenir.",
        "• AI Analiz Süreci: 360° RAG veritabanı taraması ile proaktif DÖF planı hazırlanır.",
        "• Mühendis Onayı: Sistem aksiyon emrini otomatik olarak doğrudan sahaya göndermez.",
        "• Sahaya Sevk: Başmühendis kararı inceleyip '📌 Onayla ve Sahaya Sevk Et' butonuna bastığı an canlı görev panosuna aktarılır."
    ]

    for pt in points:
        p_pt = tf_wf.add_paragraph()
        p_pt.text = f"{pt}\n"
        p_pt.font.size = Pt(16)
        p_pt.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLAYT 6: Gelecek Vizyonu & Kapanış
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = DARK_BG
    bg6.line.color.rgb = DARK_BG

    tb_end = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.0))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True

    p_e1 = tf_end.paragraphs[0]
    p_e1.text = "🎯 Vizyon: Sıfır Kaza Hedefi"
    p_e1.font.size = Pt(36)
    p_e1.font.bold = True
    p_e1.font.color.rgb = ORANGE
    p_e1.alignment = PP_ALIGN.CENTER

    p_e2 = tf_end.add_paragraph()
    p_e2.text = "\n\"Arın AI: Madende Güvenlik Tesadüf Değil, Mühendisliktir.\"\n"
    p_e2.font.size = Pt(22)
    p_e2.font.italic = True
    p_e2.font.color.rgb = WHITE
    p_e2.alignment = PP_ALIGN.CENTER

    p_e3 = tf_end.add_paragraph()
    p_e3.text = "Teşekkürler! / Sorular & Canlı Demo"
    p_e3.font.size = Pt(18)
    p_e3.font.color.rgb = LIGHT_GRAY
    p_e3.alignment = PP_ALIGN.CENTER

    prs.save("Arin_AI_Sunum.pptx")
    print("✅ 'Arin_AI_Sunum.pptx' başarıyla oluşturuldu!")

if __name__ == "__main__":
    create_arin_ai_presentation()